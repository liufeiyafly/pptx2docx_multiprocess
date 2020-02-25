# -*- coding: utf-8 -*-
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn
import time

from multiprocessing import Pool
import os

def Group_digui(group, l=[]):  # 递归寻找group组合对象里面的文字，（group里面可能套了group）
    auto_object = '<class \'pptx.shapes.autoshape.Shape\'>'
    place_object = '<class \'pptx.shapes.placeholder.SlidePlaceholder\'>'
    group_object = '<class \'pptx.shapes.group.GroupShape\'>'
    for i in group.shapes:
        if str(type(i)) == auto_object or str(type(i)) == place_object:
            s2 = i.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
            l.append(s2.strip().replace('\x0b', ' '))
        elif i.shape_type == 6:  # 等于6即 是group组合
            Group_digui(i)  # 递归访问即可
    return l  # 返回一个列表，里面元素是group内所有的文字。
def ppt2pptx(filename):

    doc = Document()
    doc.styles['Normal'].font.name = u'宋体'
    doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')

    prs=Presentation(filename)
    num=len(prs.slides) # 幻灯片的页数
    auto_object='<class \'pptx.shapes.autoshape.Shape\'>'
    place_object='<class \'pptx.shapes.placeholder.SlidePlaceholder\'>'
    group_object='<class \'pptx.shapes.group.GroupShape\'>'
    L=[]

    for i in range(num):
        slide = prs.slides[i]  # 通过索引序号 访问每页幻灯片
        body_shapes = slide.shapes  # slide.shapes 就是每页幻灯片中的所有元素。

        for j in body_shapes:
            if str(type(j)) == auto_object or str(type(j)) == place_object:  # 图形或者文本框是有 .text属性的，可以直接访问
                s1 = j.text.encode('gbk', 'ignore').decode('gbk', 'ignore')  # 用这个代码解决不能打印和写入文件'gbk'特殊编码的问题了！！！
                L.append(s1.strip().replace('\x0b', ' '))

            elif str(type(j)) == group_object:
                text_l = Group_digui(j)
                L.extend(text_l)

            elif j.shape_type == 19:  # 19代表：表格
                for row in j.table.rows:
                    for cell in row.cells:
                        s3 = cell.text_frame.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
                        L.append(s3.strip().replace('\x0b', ' '))

    s = ' '.join(L)  # 以空格分开各个内容，也可以用回车
    doc.add_paragraph(s)
    doc.save('{}.docx'.format(filename[:-5]))  # 生成的docx保存到pptx所在文件夹

if __name__ == "__main__":
    path=r'D:\test' #选择的路径，会处理路径下所有的*.pptx
    path_list=[]
    for i in os.listdir(path):
        if i[-4:]=='pptx':
            path_list.append(os.path.join(path,i))

    start_time=time.time()

    pool=Pool()
    for file_path in path_list:
        pool.apply_async(ppt2pptx,args=(file_path,))
    pool.close()
    pool.join()

    print('耗时：{}s'.format(time.time()-start_time))